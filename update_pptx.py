import os
import sys
import openpyxl
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from datetime import datetime

def get_excel_data(wb, sheet_name, start_row=2, cols=(1, 2)):
    if sheet_name not in wb.sheetnames:
        return None
    ws = wb[sheet_name]
    data = []
    for row_idx in range(start_row, ws.max_row + 1):
        row_vals = [ws.cell(row=row_idx, column=c).value for c in cols]
        if not any(v is not None for v in row_vals):
            break
        data.append(row_vals)
    return data

def update_report(excel_file, pptx_template, output_file, month_str):
    if not os.path.exists(excel_file):
        print(f"Excel file {excel_file} not found.")
        return
    
    wb = openpyxl.load_workbook(excel_file, data_only=True)
    prs = Presentation(pptx_template)
    
    # 1. Slide 3: User Funnel & Engagement
    slide3 = prs.slides[2]
    
    # Chart 2: User Funnel (Slide 3)
    funnel_data = get_excel_data(wb, 'User_funnel', start_row=2, cols=(1, 3)) # March is Col 3
    if funnel_data:
        chart_data = CategoryChartData()
        chart_data.categories = [d[0] for d in funnel_data]
        chart_data.add_series('March', [d[1] for d in funnel_data])
        for shape in slide3.shapes:
            if shape.name == "Chart 2":
                shape.chart.replace_data(chart_data)
                
    # Chart 15: User Engagement (Slide 3)
    # Get Jan, Feb, March data if possible. User Engagement Vertical rows 67-94 (Feb)
    # For now, let's just grab the whole range used in the Excel graph
    ws_ue = wb['User Engagement']
    # Rows for Feb+Mar in Aji_game copy_March.xlsx are roughly 67 to 125?
    # Let's dynamically find Mar and Mar-1
    target_months = [month_str.lower()[:3]]
    # (Simplified for this script)
    ue_data = []
    for r in range(2, ws_ue.max_row + 1):
        date_val = ws_ue.cell(row=r, column=1).value
        if isinstance(date_val, datetime):
            m = date_val.strftime('%b').lower()
            # If user asks March, we need Feb and March
            if m in ['feb', 'mar']:
                ue_data.append((date_val.strftime('%d/%m'), ws_ue.cell(row=r, column=2).value, ws_ue.cell(row=r, column=3).value))
    
    if ue_data:
        chart_data = CategoryChartData()
        chart_data.categories = [d[0] for d in ue_data]
        chart_data.add_series('New user', [d[1] for d in ue_data])
        chart_data.add_series('returning User', [d[2] for d in ue_data])
        for shape in slide3.shapes:
            if shape.name == "Chart 15":
                shape.chart.replace_data(chart_data)

    # 2. Slide 4: Gameplay Report Score
    slide4 = prs.slides[3]
    ws_score = wb['gameplay_report(score) ']
    score_chart_data = []
    for r in range(2, ws_score.max_row + 1):
        date_val = ws_score.cell(row=r, column=1).value
        if isinstance(date_val, datetime) and date_val.strftime('%b').lower() in ['mar']:
             score_chart_data.append((date_val.strftime('%d/%m'), ws_score.cell(row=r, column=2).value))
             
    if score_chart_data:
        chart_data = CategoryChartData()
        chart_data.categories = [d[0] for d in score_chart_data]
        chart_data.add_series('Score', [d[1] for d in score_chart_data])
        for shape in slide4.shapes:
            if shape.has_chart and shape.name == "Chart 2":
                shape.chart.replace_data(chart_data)

    # 3. Slide 5: Gameplay Report Time
    slide5 = prs.slides[4]
    ws_time = wb['gameplay_report(time) ']
    time_chart_data = []
    for r in range(2, ws_time.max_row + 1):
        date_val = ws_time.cell(row=r, column=1).value
        if isinstance(date_val, datetime) and date_val.strftime('%b').lower() in ['feb', 'mar']:
             time_chart_data.append((date_val.strftime('%d/%m'), ws_time.cell(row=r, column=2).value))
             
    if time_chart_data:
        chart_data = CategoryChartData()
        chart_data.categories = [d[0] for d in time_chart_data]
        chart_data.add_series('Time', [d[1] for d in time_chart_data])
        for shape in slide5.shapes:
            if shape.has_chart and (shape.name == "Chart 10" or shape.name == "Chart 1"):
                shape.chart.replace_data(chart_data)

    # 4. Slide 6: State & Menu
    slide6 = prs.slides[5]
    state_data = get_excel_data(wb, 'state', start_row=12, cols=(1, 2))
    if state_data:
        chart_data = CategoryChartData()
        chart_data.categories = [d[0] for d in state_data]
        chart_data.add_series('Count', [d[1] for d in state_data])
        for shape in slide6.shapes:
            if shape.has_chart and shape.name == "Chart 2":
                shape.chart.replace_data(chart_data)
                
    menu_data = get_excel_data(wb, 'menu', start_row=2, cols=(1, 2))
    if menu_data:
        # Top 10 menus?
        menu_data.sort(key=lambda x: x[1] if x[1] else 0, reverse=True)
        top_menus = menu_data[:10]
        chart_data = CategoryChartData()
        chart_data.categories = [d[0] for d in top_menus]
        chart_data.add_series('Count', [d[1] for d in top_menus])
        for shape in slide6.shapes:
            if shape.has_chart and shape.name == "Chart 7":
                shape.chart.replace_data(chart_data)

    # 5. Slide 7: Leaderboard
    slide7 = prs.slides[6]
    leader_data = get_excel_data(wb, 'score', start_row=2, cols=(2, 3))
    if leader_data:
        leader_data.sort(key=lambda x: x[1] if x[1] else 0, reverse=True)
        top_leaders = leader_data[:10]
        chart_data = CategoryChartData()
        chart_data.categories = [d[0] for d in top_leaders]
        chart_data.add_series('Total Score', [d[1] for d in top_leaders])
        for shape in slide7.shapes:
            if shape.has_chart and shape.name == "Chart 1":
                shape.chart.replace_data(chart_data)

    # Text Updates (Dates etc)
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                if "2025/11/28 – 2026/03/31" in shape.text:
                    # Could update to actual range
                    pass
                if "March-2026" in shape.text and month_str == "Feb":
                     shape.text = shape.text.replace("March-2026", "February-2026")

    prs.save(output_file)
    print(f"Report saved to {output_file}")

if __name__ == "__main__":
    update_report("Aji_game copy_March.xlsx", "Merkle Thailand -Ajipanda's Kitchen report- 260331  copy.pptx", "Report_March.pptx", "March")
