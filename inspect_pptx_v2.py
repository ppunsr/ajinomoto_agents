from pptx import Presentation
import os

def inspect_pptx(file_path):
    if not os.path.exists(file_path):
        print(f"File {file_path} not found.")
        return

    prs = Presentation(file_path)
    print(f"Presentation has {len(prs.slides)} slides.")
    
    for i, slide in enumerate(prs.slides):
        print(f"\n--- Slide {i+1} ---")
        for shape in slide.shapes:
            # Check for standard charts
            if shape.has_chart:
                chart = shape.chart
                # Check for external data source
                is_external = False
                try:
                    # Accessing internal relationship to check for external link
                    rId = chart._chartSpace.xpath('//c:externalData/@r:id')
                    if rId:
                        is_external = True
                except:
                    pass
                print(f"  Chart Shape: {shape.name}, Type: {chart.chart_type}, External: {is_external}")
            
            # Check for extended charts (ChartEx)
            elif shape.shape_type == 14: # GraphicFrame
                print(f"  GraphicFrame Shape: {shape.name}")
                # We can't easily see if it's a ChartEx via python-pptx easily
                # but we can check the XML if we really wanted to.
            
            elif shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text:
                    print(f"  Text Shape: {shape.name}, Text Sample: {text[:50]}...")

if __name__ == "__main__":
    inspect_pptx("Merkle Thailand -Ajipanda's Kitchen report- 260331  copy.pptx")
