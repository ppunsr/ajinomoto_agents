from pptx import Presentation
import os

def inspect_pptx(file_path):
    if not os.path.exists(file_path):
        print(f"File {file_path} not found.")
        return

    prs = Presentation(file_path)
    print(f"Presentation has {len(prs.slides)} slides.")
    
    for i, slide in enumerate(prs.slides):
        print(f"\n--- Slide {i+1} ---")
        for shape in slide.shapes:
            if shape.has_chart:
                chart = shape.chart
                print(f"  Chart Shape: {shape.name}, Type: {chart.chart_type}")
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text:
                    print(f"  Text Shape: {shape.name}, Text Sample: {text[:100]}...")

if __name__ == "__main__":
    inspect_pptx("Merkle Thailand -Ajipanda's Kitchen report- 260331  copy.pptx")
